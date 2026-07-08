"""Loaders para los distintos formatos de documentos de RutaSur."""
import json
from pathlib import Path

from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown as md_lib


def load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    texto = []
    for i, page in enumerate(reader.pages, 1):
        contenido = page.extract_text() or ""
        if contenido.strip():
            texto.append(f"[Página {i}]\n{contenido}")
    return "\n\n".join(texto)


def load_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    partes = []
    for p in doc.paragraphs:
        if p.text.strip():
            partes.append(p.text)
    # También extraer tablas
    for tabla in doc.tables:
        for fila in tabla.rows:
            celdas = [c.text.strip() for c in fila.cells if c.text.strip()]
            if celdas:
                partes.append(" | ".join(celdas))
    return "\n".join(partes)


def load_xlsx(path: Path) -> str:
    wb = load_workbook(str(path), data_only=True)
    partes = []
    for hoja in wb.worksheets:
        partes.append(f"[Hoja: {hoja.title}]")
        for fila in hoja.iter_rows(values_only=True):
            valores = [str(v) for v in fila if v is not None]
            if valores:
                partes.append(" | ".join(valores))
    return "\n".join(partes)


def load_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    partes = []
    for i, slide in enumerate(prs.slides, 1):
        partes.append(f"[Diapositiva {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texto = "".join(run.text for run in para.runs)
                    if texto.strip():
                        partes.append(texto)
    return "\n".join(partes)


def load_json(path: Path) -> str:
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    # Convertir JSON a texto legible
    return json.dumps(data, ensure_ascii=False, indent=2)


def load_html(path: Path) -> str:
    with open(path, encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def load_markdown(path: Path) -> str:
    with open(path, encoding="utf-8") as f:
        texto_md = f.read()
    html = md_lib.markdown(texto_md, extensions=["tables"])
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n", strip=True)


def load_csv(path: Path) -> str:
    with open(path, encoding="utf-8") as f:
        return f.read()


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
    ".json": load_json,
    ".html": load_html,
    ".md": load_markdown,
    ".csv": load_csv,
}


def load_document(path: Path) -> str:
    """Detecta el formato por extensión y extrae el texto."""
    ext = path.suffix.lower()
    loader = LOADERS.get(ext)
    if not loader:
        raise ValueError(f"Formato no soportado: {ext}")
    return loader(path)


def load_all_documents(carpeta: str = "./documents") -> dict[str, str]:
    """Carga todos los documentos de la carpeta. Devuelve {nombre: texto}."""
    resultados = {}
    for path in sorted(Path(carpeta).iterdir()):
        if path.suffix.lower() in LOADERS:
            try:
                texto = load_document(path)
                resultados[path.name] = texto
                print(f"✅ {path.name}: {len(texto)} caracteres")
            except Exception as e:
                print(f"❌ {path.name}: {e}")
    return resultados
